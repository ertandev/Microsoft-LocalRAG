import os
import sqlite3
import json
import csv
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from foundry_local_sdk import Configuration, FoundryLocalManager

# 1. Dosya Türlerine Göre Metin Okuma Fonksiyonları

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text_runs = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_runs.append(text)
    return "\n".join(text_runs)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text_runs = []
    for p in doc.paragraphs:
        if p.text:
            text_runs.append(p.text)
    # Tabloları da oku
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells if cell.text]
            if row_text:
                text_runs.append(" | ".join(row_text))
    return "\n".join(text_runs)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_xlsx(file_path):
    wb = load_workbook(file_path, read_only=True)
    text_runs = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text_runs.append(" | ".join(row_text))
    return "\n".join(text_runs)

def extract_text_from_csv(file_path):
    text_runs = []
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            if row:
                text_runs.append(" , ".join(row))
    return "\n".join(text_runs)

def extract_text_from_json(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)

def extract_text_from_html(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
        # Script ve style etiketlerini temizle
        for s in soup(["script", "style"]):
            s.decompose()
        return soup.get_text(separator='\n')

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext in ['.txt', '.md']:
        return extract_text_from_txt(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext == '.json':
        return extract_text_from_json(file_path)
    elif ext in ['.html', '.htm']:
        return extract_text_from_html(file_path)
    else:
        print(f"Uyarı: '{ext}' dosya türü desteklenmiyor, atlanıyor: {file_path}")
        return None

def get_chunks(text, chunk_size=600, overlap=120):
    if not text:
        return []
    
    import re
    # Simple split by sentence endings
    sentences = re.split(r'(?<=[.?!])\s+', text.strip())
    
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        
        # If a single sentence exceeds the chunk_size, split by words
        if len(sentence) > chunk_size:
            if current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_length = 0
            
            words = sentence.split(' ')
            sub_chunk = []
            sub_len = 0
            for w in words:
                if sub_len + len(w) + 1 > chunk_size:
                    chunks.append(" ".join(sub_chunk))
                    # Overlap with last 5 words
                    overlap_words = sub_chunk[-5:] if len(sub_chunk) >= 5 else sub_chunk
                    sub_chunk = list(overlap_words) + [w]
                    sub_len = sum(len(x) + 1 for x in sub_chunk)
                else:
                    sub_chunk.append(w)
                    sub_len += len(w) + 1
            if sub_chunk:
                current_chunk = sub_chunk
                current_length = sub_len
            continue

        if current_length + len(sentence) + 1 > chunk_size:
            if current_chunk:
                chunks.append(" ".join(current_chunk))
            
            # Form overlap: take last sentences that fit within overlap limit
            overlap_chunk = []
            overlap_len = 0
            for s in reversed(current_chunk):
                if overlap_len + len(s) + 1 <= overlap:
                    overlap_chunk.insert(0, s)
                    overlap_len += len(s) + 1
                else:
                    break
            
            current_chunk = overlap_chunk + [sentence]
            current_length = sum(len(x) + 1 for x in current_chunk)
        else:
            current_chunk.append(sentence)
            current_length += len(sentence) + 1
            
    if current_chunk:
        chunks.append(" ".join(current_chunk))
        
    return [c.strip() for c in chunks if c.strip()]

# 3. Ana Çalıştırma Bloğu
def main():
    db_path = os.path.join(os.path.dirname(__file__), "knowledge_base.db")
    docs_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "documents")

    # Eğer documents klasörü yoksa oluştur
    if not os.path.exists(docs_dir):
        os.makedirs(docs_dir)
        print(f"'{docs_dir}' klasörü oluşturuldu. Lütfen içine okutmak istediğiniz belgeleri atın.")
        return

    # Klasördeki dosyaları kontrol et
    files = [os.path.join(docs_dir, f) for f in os.listdir(docs_dir) if os.path.isfile(os.path.join(docs_dir, f))]
    if not files:
        print(f"'{docs_dir}' klasörü boş. Lütfen içine pdf, docx, txt vb. dökümanlar yerleştirin.")
        return

    print("Foundry Local başlatılıyor...")
    config = Configuration(app_name="local_rag_app")
    try:
        FoundryLocalManager.initialize(config)
    except Exception:
        pass
    manager = FoundryLocalManager.instance

    print("Embedding modeli yükleniyor...")
    embed_model = manager.catalog.get_model("qwen3-embedding-8b")
    embed_model.load()
    embedding_client = embed_model.get_embedding_client()

    # Veritabanı tablosunu hazırla
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            text TEXT,
            embedding TEXT,
            file_name TEXT
        )
    """)
    conn.commit()

    # Her çalıştırmada eski dökümanları temizleyelim (çift kayıt olmaması için)
    cursor.execute("DELETE FROM documents")
    conn.commit()

    print("\n--- BELGELER DİZİNLENİYOR ---")
    for file_path in files:
        file_name = os.path.basename(file_path)
        print(f"\nDosya okunuyor: {file_name}")
        
        # Metni dosyadan çıkar
        raw_text = extract_text(file_path)
        if not raw_text:
            continue

        # Metni küçük anlamlı parçalara ayır (chunking)
        chunks = get_chunks(raw_text)
        print(f"-> Dosya {len(chunks)} parçaya bölündü.")

        # Her parçayı sayısal vektöre çevirip SQLite'a kaydet
        for i, chunk in enumerate(chunks):
            if not chunk.strip():
                continue
            
            response = embedding_client.generate_embedding(chunk)
            vector = response.data[0].embedding
            vector_str = json.dumps(vector)

            cursor.execute("INSERT INTO documents (text, embedding, file_name) VALUES (?, ?, ?)", (chunk, vector_str, file_name))
            print(f"   [Parça {i+1}/{len(chunks)}] Kaydedildi. Vektör boyutu: {len(vector)}")

    conn.commit()
    conn.close()
    print("\n==========================================")
    print("DİZİNLEME TAMAMLANDI! Tüm belgeler veri tabanına işlendi.")
    print("React uygulamanızı yenileyerek belgeleri görebilirsiniz.")
    print("==========================================\n")

if __name__ == "__main__":
    main()
